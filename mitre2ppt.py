import re
from bs4 import BeautifulSoup
import requests
from pptx import Presentation
from pptx.util import Pt


def get_mitre(url):
    # define client and prepare header
    user_agent = 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/70.0.3538.102 Safari/537.36'
    headers = {'User-Agent': user_agent}
    try: 
        # send get request to MITRE
        response = requests.get(url, headers=headers)
    except Exception as e:
        print('Error accessing webpage. Error Message: ', str(e))
    # print the response status code to the console
    print('Response: ', response)
    return response


def main():
    response = get_mitre('https://attack.mitre.org/techniques/enterprise/')
    try:
        # parse webpage
        website_content = BeautifulSoup(response.content, 'html.parser')
        # find table and its rows
        table_body = website_content.find('tbody')
        rows = table_body.find_all('tr')
        # create power point basis and the layout to be used
        ppt = Presentation()
        layout = ppt.slide_layouts[1]
        # for each row displayed in the table on the webpage
        for entry in rows:
            # get the table data and clean it from html tags
            contents = entry.find_all('td')
            description = contents[len(contents)-1].text.replace('\n', '').replace('  ', '')
            name = contents[len(contents)-2].text.replace('\n', '')
            # add slide and initialize title and content text box
            slide = ppt.slides.add_slide(layout)
            title = slide.placeholders[0]
            content_text = slide.placeholders[1]
            # set the title of the slide = name of the MITRE Technique
            if entry.attrs['class'][0] == "technique":
                technique = name
                title.text = name
                # for techniques also get the tactics from the technique page
                sub_page = entry.contents[1].contents[1].attrs['href']
                technique_page = get_mitre('https://attack.mitre.org/' + sub_page)
                technique_page = BeautifulSoup(technique_page.content, 'html.parser')
                # get tactics from webpage
                tactics = technique_page.find('div', {'id': 'card-tactics'})
                # clean extracted content
                tactics = tactics.text.replace('\n', '').replace('ⓘTactics:', '').replace('ⓘTactic:', '')
                tactics = re.sub(re.compile('[ ]{2,5}'), '', str(tactics))
            else:
                title.text = technique + "-" + name
            # set the font size for the title
            slide.placeholders[0].text_frame.paragraphs[0].font.size = Pt(32)
            # set the content of the slide = description of the MITRE Technique + additional text/info if needed
            content_text.text = description + '\n\nTactics:\t\t' + tactics + '\nRelevant?\t\tYes/No'
            # set the font size of the description for each paragraph
            i = len(slide.placeholders[1].text_frame.paragraphs) - 1
            while i >= 0:
                slide.placeholders[1].text_frame.paragraphs[i].font.size = Pt(16)
                i -= 1
        # save the power point file
        ppt.save('MITRE ATT&CK.pptx')
    except Exception as e:
        print('Error creating power point file. Error Message: ', str(e))
    print('Finished creating MITRE pptx file!')


if __name__ == '__main__':
    main()